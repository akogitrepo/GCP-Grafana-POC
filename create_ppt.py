#!/usr/bin/env python3
"""
GCP Platform Leadership Review — PowerPoint deck.
All headers use ONLY colored rectangles + text shapes (zero images).
This eliminates every image-sizing / overlap bug seen in previous versions.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width   # shorthand for full slide width

# ── Brand colours ─────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x0A, 0x0F, 0x3D)   # very dark navy  → branding bar bg
C_DBLUE  = RGBColor(0x1A, 0x23, 0x7E)   # dark blue       → title bars / footer
C_GBLUE  = RGBColor(0x42, 0x85, 0xF4)   # Google blue
C_GREEN  = RGBColor(0x34, 0xA8, 0x53)   # Google green
C_ORANGE = RGBColor(0xFB, 0xBC, 0x04)   # Google yellow
C_RED    = RGBColor(0xEA, 0x43, 0x35)   # Google red
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DGRAY  = RGBColor(0x33, 0x33, 0x33)   # body text
C_LBLUE  = RGBColor(0x7B, 0xB3, 0xF0)   # light-blue text on dark bg
C_VLBLUE = RGBColor(0xBB, 0xDE, 0xFB)   # very-light blue  (subtitle text)
C_MIDBL  = RGBColor(0x0F, 0x17, 0x5E)   # mid-navy for footer on title slides

GOOGLE4 = [C_GBLUE, C_RED, C_ORANGE, C_GREEN]

# ── Fixed layout measurements (content slides) ────────────────────────────────
# All values in EMU.  Content starts at CTOP — guaranteed gap below header.
BRAND_H  = Inches(0.42)          # layer 1 — dark navy branding bar
STRIP_H  = Inches(0.10)          # layer 2 — 4-colour accent strip
TITLB_H  = Inches(0.65)          # layer 3 — dark-blue title bar  (22 pt fits fine)
HDR_H    = BRAND_H + STRIP_H + TITLB_H   # 1.17" total header height
CTOP     = HDR_H + Inches(0.23)  # 1.40" — where body content starts
FTR_TOP  = Inches(7.00)          # footer top
FTR_H    = Inches(0.45)          # footer height
BODY_H   = FTR_TOP - CTOP - Inches(0.05) # ≈ 5.55" available for body

# Convenience floats (inches) for the add_body() helper
CTOP_IN  = CTOP / 914400         # ≈ 1.40
BODY_H_IN = BODY_H / 914400      # ≈ 5.55


# ═══════════════════════════════════════════════════════════════════════════════
# SHARED HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _rect(slide, left, top, width, height, fill_rgb, line=False):
    """Add a filled rectangle; line=False removes border."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    if line:
        s.line.color.rgb = fill_rgb
    else:
        s.line.fill.background()
    return s


def _textbox(slide, left_in, top_in, w_in, h_in, text, size_pt,
             color=C_DGRAY, bold=False, align=PP_ALIGN.LEFT,
             wrap=True, anchor=MSO_ANCHOR.TOP):
    """Add a simple single-paragraph text box (dimensions in inches)."""
    tx = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tx


def add_footer(slide, on_dark=False):
    """Footer bar at bottom of every slide."""
    bar = _rect(slide, 0, FTR_TOP, W, FTR_H,
                C_MIDBL if on_dark else C_DBLUE)
    tf = bar.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = Inches(0.5)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "AI Engineer / SRE Architect – Abhay Ojha"
    p.font.size = Pt(11)
    p.font.color.rgb = C_LBLUE if on_dark else C_WHITE
    p.font.bold = False
    p.alignment = PP_ALIGN.RIGHT


def add_header(slide, title_text):
    """
    Three-layer header — PURE SHAPES, no images.

    Layer 1  (0          → BRAND_H  )  very dark navy  +  GCP branding text
    Layer 2  (BRAND_H    → +STRIP_H )  4-colour Google accent strip
    Layer 3  (+STRIP_H   → +TITLB_H )  dark-blue title bar  +  slide title
    """
    # ── Layer 1: branding bar ─────────────────────────────────────────────────
    nav = _rect(slide, 0, 0, W, BRAND_H, C_NAVY)
    tf_n = nav.text_frame
    tf_n.word_wrap = False
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.margin_left   = Inches(0.5)
    tf_n.margin_right  = Inches(0.5)
    tf_n.margin_top    = tf_n.margin_bottom = 0
    p_n = tf_n.paragraphs[0]
    p_n.text = "GCP × SRE × O11y   ·   PRODUCTION-GRADE · GA"
    p_n.font.size = Pt(10)
    p_n.font.color.rgb = C_LBLUE
    p_n.font.bold = False
    p_n.alignment = PP_ALIGN.LEFT

    # ── Layer 2: 4-colour accent strip ────────────────────────────────────────
    seg_w = W // 4
    for i, c in enumerate(GOOGLE4):
        _rect(slide, seg_w * i, BRAND_H, seg_w, STRIP_H, c)

    # ── Layer 3: dark-blue title bar ──────────────────────────────────────────
    title_top = BRAND_H + STRIP_H
    bar = _rect(slide, 0, title_top, W, TITLB_H, C_DBLUE)
    tf = bar.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left   = Inches(0.5)
    tf.margin_right  = Inches(0.5)
    tf.margin_top    = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT


def add_body(slide, bullets, left_in=0.6, top_in=None,
             width_in=12.1, height_in=None, font_size=17):
    """
    Add a bulleted text box.  Defaults put it in the standard content zone.
    bullets = list of (text, indent_level) tuples.
    """
    if top_in is None:
        top_in = CTOP_IN
    if height_in is None:
        height_in = BODY_H_IN
    tx = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (text, lvl) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = C_DGRAY
        p.level = lvl
        p.space_after = Pt(5)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _table_header_row(tbl, col_data, bg=C_DBLUE):
    """Style first row of a table as a header."""
    for c, val in enumerate(col_data):
        cell = tbl.cell(0, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = C_WHITE
        p.font.bold = True


def _table_data_row(tbl, r, row_data):
    for c, val in enumerate(row_data):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = C_DGRAY


# ═══════════════════════════════════════════════════════════════════════════════
# TITLE SLIDE  (Slide 1)
# ═══════════════════════════════════════════════════════════════════════════════
# Layout
#   0.00" → 2.30"  Dark navy banner (branding area, mimics banner.png)
#   2.30" → 2.40"  4-colour accent strip
#   2.40" → 6.95"  Dark-blue content area (title, subtitle, category)
#   7.00" → 7.45"  Footer

TITLE_BANNER_H = Inches(2.30)
TITLE_STRIP_H  = Inches(0.10)
TITLE_BODY_TOP = TITLE_BANNER_H + TITLE_STRIP_H   # 2.40"

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Full-slide dark-blue background
_rect(slide1, 0, 0, W, prs.slide_height, C_DBLUE)

# ── Banner area (dark navy) ───────────────────────────────────────────────────
_rect(slide1, 0, 0, W, TITLE_BANNER_H, C_NAVY)

# "● PRODUCTION-GRADE · GA"
_textbox(slide1, 0.6, 0.22, 6, 0.35,
         "●  PRODUCTION-GRADE · GA", 11, C_LBLUE, bold=False)

# Large heading inside banner
_textbox(slide1, 0.6, 0.58, 9, 1.10,
         "GCP × SRE × O11y", 52, C_WHITE, bold=True)

# Sub-heading inside banner
_textbox(slide1, 0.6, 1.58, 11, 0.45,
         "Multi-region GKE platform with Grafana-on-BigQuery observability",
         17, C_LBLUE)

# Feature pill row (simulate with text; rounded-rect pills via shapes)
pill_labels = ["2 GKE clusters", "Multi-Cluster Ingress",
               "Cloud Armor WAF", "BigQuery + Grafana"]
pill_x = [0.55, 2.85, 5.75, 8.45]
for lbl, px in zip(pill_labels, pill_x):
    pill = _rect(slide1, Inches(px), Inches(1.90), Inches(2.10), Inches(0.30), C_DBLUE)
    # pill border
    pill.line.color.rgb = C_GBLUE
    tf_p = pill.text_frame
    tf_p.word_wrap = False
    tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_p.margin_left = tf_p.margin_right = Inches(0.08)
    tf_p.margin_top = tf_p.margin_bottom = 0
    pp = tf_p.paragraphs[0]
    pp.text = lbl
    pp.font.size = Pt(11)
    pp.font.color.rgb = C_VLBLUE
    pp.alignment = PP_ALIGN.CENTER

# 4-colour accent strip
seg_w = W // 4
for i, c in enumerate(GOOGLE4):
    _rect(slide1, seg_w * i, TITLE_BANNER_H, seg_w, TITLE_STRIP_H, c)

# ── Main title text (below strip) ────────────────────────────────────────────
_textbox(slide1, 0.8, 2.60, 11.7, 1.00,
         "GCP Multi-Cluster GKE Platform",
         42, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

_textbox(slide1, 0.8, 3.58, 11.7, 0.55,
         "Two Web Apps  ·  Global Load Balancing  ·  Grafana-on-BigQuery Observability",
         20, C_VLBLUE, align=PP_ALIGN.CENTER)

_textbox(slide1, 0.8, 4.55, 11.7, 0.45,
         "Leadership Review  |  Platform Engineering",
         17, C_LBLUE, align=PP_ALIGN.CENTER)

add_footer(slide1, on_dark=True)

add_notes(slide1, """MEETING NOTES — Title Slide
- Welcome leadership team
- Purpose: Review the GCP multi-cluster platform architecture, delivery plan, and key decisions
- Platform: two stateless web apps across two GKE clusters in different regions
- Fronted by a global HTTPS load balancer; full observability via Grafana + BigQuery
- All infrastructure is Terraform — single 'terraform apply' provisions everything""")


# ═══════════════════════════════════════════════════════════════════════════════
# CONTENT SLIDES  (Slides 2–16)
# ═══════════════════════════════════════════════════════════════════════════════

# ── Slide 2: Agenda ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Agenda")
add_footer(s)
add_body(s, [
    ("1.   Executive Summary & Business Value", 0),
    ("2.   Architecture Overview", 0),
    ("3.   GKE Cluster Design", 0),
    ("4.   Application Deployment Strategy", 0),
    ("5.   Traffic Flow & Global Load Balancing", 0),
    ("6.   Security Posture", 0),
    ("7.   Observability & Grafana Dashboards", 0),
    ("8.   High Availability & Disaster Recovery", 0),
    ("9.   CI/CD Pipeline", 0),
    ("10.  Project Timeline & Effort", 0),
    ("11.  Key Design Decisions", 0),
    ("12.  Risks & Mitigations", 0),
    ("13.  Next Steps & Ask", 0),
], font_size=19)
add_notes(s, """MEETING NOTES — Agenda
- Walk through agenda so leadership knows what to expect
- Invite questions at any point or hold for end
- Total time: ~30–40 minutes
- Key decision points will be called out on relevant slides""")


# ── Slide 3: Executive Summary ────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Executive Summary")
add_footer(s)
add_body(s, [
    ("What We're Building", 0),
    ("A production-grade, multi-region GCP platform running two web apps with global load balancing, full observability, and zero-downtime failover", 1),
    ("", 0),
    ("Why It Matters", 0),
    ("Sub-second regional failover — no DNS changes, no customer impact", 1),
    ("Single Terraform apply provisions the entire stack — repeatable, auditable, drift-free", 1),
    ("Enterprise security: Workload Identity, Binary Authorization, Cloud Armor WAF, no JSON keys", 1),
    ("Grafana dashboards on BigQuery provide deep operational visibility + forensic investigation", 1),
    ("", 0),
    ("Key Numbers", 0),
    ("2 regions (us-central1, us-east1)  ·  2 GKE clusters  ·  2 apps  ·  4 Grafana panels", 1),
    ("~15.5 engineer-days total effort  ·  7–20 calendar days depending on team size", 1),
], font_size=17)
add_notes(s, """MEETING NOTES — Executive Summary
- Complete platform-as-code: not just infra, but apps, CI/CD, observability, and security
- Multi-Cluster Ingress: instant failover (seconds, not minutes like DNS-based)
- All infra is version-controlled Terraform — no manual console clicks
- Security baked in: no SA keys, WAF at the edge, binary authorization blocks unattested images
- 15.5 engineer-day estimate covers everything from VPC to Grafana dashboards
- With 2 engineers: 10–12 calendar days""")


# ── Slide 4: Architecture Overview ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Architecture Overview")
add_footer(s)
ARCH_IMG = "/Users/abojha/Anthropic-POC/GCP-Grafana/GCP - Grafana/docs/architecture-diagram.png"
# Image fills the content zone only — header is already above it
s.shapes.add_picture(
    ARCH_IMG,
    Inches(0.35), CTOP,
    Inches(12.6), BODY_H - Inches(0.1))
add_notes(s, """MEETING NOTES — Architecture Overview
- Diagram: customer edge → GCP project → 2 regional GKE clusters → observability
- Customer traffic top-to-bottom; observability pipeline in amber on right
- Two identical regional clusters (us-central1 primary, us-east1 secondary/DR)
- Each cluster: system + apps node pools, web-app-a and web-app-b namespaces
- Terraform module chain: project → network → iam → gke-cluster (×2) → mci → cloud-armor → observability""")


# ── Slide 5: GKE Cluster Design ───────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "GKE Cluster Design")
add_footer(s)

rows, cols = 9, 3
tbl_shape = s.shapes.add_table(rows, cols,
    Inches(0.7), CTOP, Inches(11.9), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.6)
tbl.columns[1].width = Inches(4.1)
tbl.columns[2].width = Inches(4.1)

table_data = [
    ["Setting",            "Primary (us-central1)",      "Secondary (us-east1)"],
    ["Mode",               "GKE Standard (regional)",    "GKE Standard (regional)"],
    ["Networking",         "VPC-native, private",        "VPC-native, private"],
    ["Workload Identity",  "Enabled",                    "Enabled"],
    ["Managed Prometheus", "Enabled",                    "Enabled"],
    ["Release Channel",    "REGULAR",                    "REGULAR"],
    ["Binary Authorization","Enforce-policy",            "Enforce-policy"],
    ["Node Pools",         "system, apps",               "system, apps"],
    ["Shielded Nodes",     "Enabled",                    "Enabled"],
]
_table_header_row(tbl, table_data[0])
for r in range(1, rows):
    _table_data_row(tbl, r, table_data[r])

add_notes(s, """MEETING NOTES — GKE Cluster Design
- Both clusters identical by design — symmetric failure modes
- Regional Standard over Autopilot: node-pool tuning, cost control at scale
- Workload Identity eliminates SA key management
- Managed Prometheus on both clusters for metric collection
- Binary Authorization enforce mode blocks unattested container images
- Separate system/apps node pools prevent workload interference""")


# ── Slide 6: Application Deployment ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Application Deployment Strategy")
add_footer(s)
add_body(s, [
    ("Two Stateless Web Applications", 0),
    ("web-app-a: Purely stateless (Cloud SQL HA optional)", 1),
    ("web-app-b: Adds Pub/Sub + Memorystore — included in /readyz health checks", 1),
    ("", 0),
    ("Deployment Model (identical for both apps)", 0),
    ("Multi-pod Deployment with Kustomize base + per-region overlays (primary/secondary)", 1),
    ("ConfigMap-driven configuration (region, log level, app name)", 1),
    ("Workload Identity for downstream auth — no baked-in secrets", 1),
    ("HPA on CPU + memory for auto-scaling", 1),
    ("PDB (minAvailable: 2) for zero-downtime rollouts", 1),
    ("Topology spread across zones for real availability", 1),
    ("Dependency-aware /readyz — Redis outage removes pod from rotation without restart", 1),
    ("Graceful drain on SIGTERM for clean shutdown", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — Application Deployment
- Apps follow 12-factor principles: stateless, config-driven, health-aware
- Kustomize over Helm: 2 apps, no charts to publish, no parameter explosion
- PDB ensures at least 2 pods always available during voluntary disruptions
- Topology spread prevents all replicas landing on one node/zone
- web-app-b readyz includes Redis check: Memorystore outage removes pod gracefully
- Secrets go to Secret Manager, read via Workload Identity at startup""")


# ── Slide 7: Traffic Flow ─────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "End-to-End Traffic Flow")
add_footer(s)
add_body(s, [
    ("1.  DNS → Customer hits the domain; Cloud DNS returns the global anycast IP", 0),
    ("2.  TLS Termination → Google-managed SSL cert at the edge", 0),
    ("3.  Cloud Armor → OWASP WAF rules + adaptive L7 DDoS + per-IP rate limiting", 0),
    ("4.  Multi-Cluster Ingress → Routes to nearest healthy cluster via NEGs", 0),
    ("5.  Container-Native LB → Reaches healthy pods directly (bypasses kube-proxy)", 0),
    ("6.  Pod Processing → Structured JSON logs, Prometheus metrics, OTel traces", 0),
    ("7.  Response → Service → NEG → LB → Customer", 0),
    ("", 0),
    ("Path Routing", 0),
    ("/a/* → web-app-a  |  /b/* → web-app-b  |  /* → web-app-a (catch-all)", 1),
    ("", 0),
    ("Failover", 0),
    ("If primary cluster fails health checks, MCI shifts traffic to secondary — no DNS changes, sub-second", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — Traffic Flow
- Single global anycast IP: one DNS record, one cert, one entry point worldwide
- Cloud Armor evaluates BEFORE traffic hits clusters — cheaper and faster than in-cluster WAF
- MCI is the key differentiator vs DNS-based failover: no TTL delays, no client caching issues
- Container-native NEG LB skips kube-proxy — lower latency, better health checking
- All telemetry (logs, metrics, traces) emitted at pod level flows to observability stack
- Failover is automatic and transparent — NOT manual""")


# ── Slide 8: Security ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Security Posture")
add_footer(s)
add_body(s, [
    ("Identity & Access", 0),
    ("Workload Identity: KSA → GSA binding, no JSON keys in pods", 1),
    ("Persona-based IAM: Dev, Ops, SRE, CI/CD roles defined in Terraform", 1),
    ("Workload Identity Federation for CI/CD — ephemeral, repo-scoped access", 1),
    ("", 0),
    ("Network & Edge Security", 0),
    ("Private nodes with IAP-SSH for break-glass only", 1),
    ("Cloud Armor: OWASP preconfigured rules + adaptive protection", 1),
    ("Default-deny NetworkPolicy per namespace; explicit allows only", 1),
    ("", 0),
    ("Supply Chain & Data", 0),
    ("Binary Authorization in enforce mode with Cloud Build attestor", 1),
    ("Secret Manager for all secrets (mounted via CSI driver)", 1),
    ("KMS-encrypted persistent disks; Artifact Registry with Trivy scanning", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — Security
- ZERO service account JSON keys anywhere in the system
- WIF: GitHub Actions gets short-lived tokens scoped to this repo only
- Binary Authorization blocks ANY container not built by our pipeline
- Cloud Armor at LB edge: malicious traffic never reaches clusters
- NetworkPolicies are default-deny: only explicitly allowed traffic flows
- Secret Manager + CSI: secrets never in K8s objects or environment variables
- Aligns with OWASP Top 10 and GCP security best practices""")


# ── Slide 9: Observability ────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Observability & Grafana Dashboards")
add_footer(s)
add_body(s, [
    ("Telemetry Stack", 0),
    ("Container logs + K8s events → Cloud Logging (default GKE agent)", 1),
    ("Metrics → Cloud Monitoring + Managed Prometheus + kube-state-metrics", 1),
    ("Traces → Cloud Trace (OTel spans from apps)", 1),
    ("Profiles → Cloud Profiler  |  Errors → Error Reporting", 1),
    ("", 0),
    ("Grafana Dashboard — 4 Panels on BigQuery", 0),
    ("Panel 1: Application errors per minute (per service)", 1),
    ("Panel 2: Pod restarts by namespace", 1),
    ("Panel 3: Request latency p50 / p95 / p99", 1),
    ("Panel 4: Container CPU & memory utilization", 1),
    ("", 0),
    ("Data Pipeline", 0),
    ("Cloud Logging → BigQuery sink (platform_logs) → SQL Views → Grafana", 1),
    ("Partitioned by timestamp, 60-day retention, pre-aggregated views for cost control", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — Observability
- Full stack: logs, metrics, traces, profiles, error reporting
- 4 Grafana panels backed by BigQuery views — same data for ad-hoc forensic queries
- BigQuery: zero infra to run, dual-use for dashboards AND investigations
- Tradeoff: query latency is seconds not milliseconds — for sub-second metrics use Managed Prometheus
- Cost control: timestamp partitioning, 60-day expiry, filtered severity, pre-aggregated views
- Dashboard JSON is version-controlled and provisioned automatically""")


# ── Slide 10: HA & DR ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "High Availability & Disaster Recovery")
add_footer(s)
add_body(s, [
    ("Multi-Region Failover", 0),
    ("Two regional clusters — MCI auto-fails over within seconds when health checks drop", 1),
    ("No DNS changes required — global anycast IP stays the same", 1),
    ("", 0),
    ("Data Resilience", 0),
    ("Cloud SQL HA + cross-region read replicas for persisted state", 1),
    ("Memorystore Standard tier with HA + cross-region async replication", 1),
    ("Artifact Registry is multi-regional (us) — image pull works during regional outage", 1),
    ("", 0),
    ("Infrastructure State", 0),
    ("Terraform state in GCS with versioning + 7-day soft-delete", 1),
    ("GKE control plane etcd backed up by managed control plane (no operator action)", 1),
    ("", 0),
    ("Recovery Testing", 0),
    ("DR drill included in sprint plan: kill primary cluster pods, confirm MCI failover", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — HA & DR
- MCI failover is sub-second and requires ZERO operator intervention
- Fundamentally better than DNS-based failover (TTL often 5–30 minutes)
- Cloud SQL HA: synchronous within-region; cross-region replicas are async
- Artifact Registry multi-regional: us-central1 outage doesn't block image pulls in us-east1
- Terraform state protected by GCS versioning — state corruption is recoverable
- DR drill on Day 9 — we prove failover works before handover
- Stateless apps have RPO=0, RTO=seconds""")


# ── Slide 11: CI/CD ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "CI/CD Pipeline")
add_footer(s)
add_body(s, [
    ("Dual Pipeline Support", 0),
    ("Cloud Build: Build → Trivy scan → Push to Artifact Registry → Kustomize deploy to both clusters", 1),
    ("GitHub Actions: Same flow, authenticating via Workload Identity Federation (no JSON keys)", 1),
    ("", 0),
    ("Authentication", 0),
    ("GitHub OIDC + Workload Identity Federation → ephemeral, repo-scoped access", 1),
    ("CI/CD service account with least-privilege IAM bindings", 1),
    ("", 0),
    ("Supply Chain Security", 0),
    ("Trivy vulnerability scanning on every build", 1),
    ("Binary Authorization attestor signs images built by the pipeline", 1),
    ("Unsigned/unattested images are rejected at cluster admission", 1),
    ("", 0),
    ("Deployment Strategy", 0),
    ("Deploy-on-merge to both clusters via Kustomize overlays", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — CI/CD
- Two pipeline options: Cloud Build (GCP-native) and GitHub Actions (developer-familiar)
- Both: build → scan → push → deploy
- WIF: NO service account keys stored in GitHub secrets
- Trivy catches known CVEs before images reach the registry
- Binary Authorization last gate: even direct registry pushes are rejected by the cluster
- bootstrap.sh sets up the WIF pool and provider in one command
- Terraform changes go through separate workflow: plan → approve → apply""")


# ── Slide 12: Timeline ────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Project Timeline & Effort")
add_footer(s)

rows, cols = 8, 3
tbl_shape = s.shapes.add_table(rows, cols,
    Inches(0.6), CTOP, Inches(12.1), Inches(3.7))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.1)
tbl.columns[1].width = Inches(8.4)
tbl.columns[2].width = Inches(2.6)

timeline_data = [
    ["#",    "Workstream",                                              "Eng-Days"],
    ["1–2",  "Project / IAM / VPC / NAT / GKE Clusters (×2)",          "3.5"],
    ["3–4",  "Multi-Cluster Ingress + Cloud Armor",                     "2.0"],
    ["5–6",  "Artifact Registry + Binary Auth + App Skeletons",         "1.5"],
    ["7–8",  "Kubernetes Manifests (Kustomize) + CI/CD",                "3.0"],
    ["9–10", "Observability: Log Sink → BigQuery → Grafana",            "2.0"],
    ["11–14","Alerts, Hardening, Testing, Documentation",               "3.5"],
    ["",     "TOTAL",                                                   "~15.5 days"],
]
_table_header_row(tbl, timeline_data[0])
for r in range(1, 7):
    _table_data_row(tbl, r, timeline_data[r])
# Total row — accent colour
for c, val in enumerate(timeline_data[7]):
    cell = tbl.cell(7, c)
    cell.text = val
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_GBLUE
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(13)
    p.font.color.rgb = C_WHITE
    p.font.bold = True

# Calendar note below table
note_top = CTOP_IN + 3.9
add_body(s, [
    ("Calendar Timeline", 0),
    ("Solo: 15–20 days  |  2 engineers: 10–12 days  |  3+ engineers: 7–9 days", 1),
    ("Add ~3 days buffer for org / billing / DNS provisioning waits", 1),
], top_in=note_top, height_in=1.4, font_size=16)

add_notes(s, """MEETING NOTES — Timeline
- 15.5 engineer-days is TOTAL effort across all workstreams
- 2-person team fits in a single 2-week sprint
- Week 1: Infrastructure (VPC, GKE, MCI, apps deployed end-to-end)
- Week 2: CI/CD, observability, hardening, DR drill, documentation
- 3-day buffer covers real-world waits: org policy, billing, DNS propagation, SSL cert
- All estimates assume engineers with prior GCP + GKE + Terraform experience""")


# ── Slide 13: Key Design Decisions ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Key Design Decisions")
add_footer(s)
add_body(s, [
    ("GKE Standard over Autopilot — node-pool tuning flexibility, cost control at scale", 0),
    ("Multi-Cluster Ingress over DNS-based failover — instant failover, single IP/cert", 0),
    ("BigQuery as Grafana datasource — zero infra, dual-use (dashboards + forensics)", 0),
    ("Kustomize over Helm — simpler for 2 apps, no chart publishing overhead", 0),
    ("Workload Identity Federation for CI/CD — ephemeral tokens, no key rotation burden", 0),
    ("Cloud Armor at LB edge — inspect before traffic hits cluster, lower cost", 0),
    ("Binary Authorization enforce mode — supply chain security from day one", 0),
    ("VPC-native (alias IPs) — required by MCI, clean Pod/Service CIDRs", 0),
    ("Topology spread + PDB + readyz — actual availability, not just replica count", 0),
], font_size=17)
add_notes(s, """MEETING NOTES — Key Design Decisions
- Each decision is documented with rationale in docs/design-decisions.md
- Standard vs Autopilot: Autopilot constrains daemonsets, node taints, kernel tuning
- MCI vs DNS: DNS is TTL-bound (minutes); MCI is health-check-bound (seconds)
- BigQuery tradeoff: query latency seconds not milliseconds — acceptable for 4 panels
- Kustomize: base + overlays maps cleanly to 2-app, 2-region topology
- Every decision is designed to be reversible with localized changes""")


# ── Slide 14: Risks & Mitigations ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Risks & Mitigations")
add_footer(s)

rows, cols = 6, 2
tbl_shape = s.shapes.add_table(rows, cols,
    Inches(0.6), CTOP, Inches(12.1), Inches(4.5))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(5.2)
tbl.columns[1].width = Inches(6.9)

risks = [
    ["Risk", "Mitigation"],
    ["Managed SSL cert stuck in PROVISIONING",
     "Confirm DNS A-record propagation before declaring done"],
    ["MCI features disabled at project level",
     "gke-hub + multiclusteringress APIs enabled in Terraform module"],
    ["BinAuth enforce blocks first deploy",
     "Run dry-run mode for first 24 h, then flip to enforce"],
    ["BigQuery sink hits quota on log volume",
     "Filter aggressively (severity≥INFO), partitioning ON, 60-day expiry"],
    ["Region-pair latency impacts UX",
     "Add 3rd cluster (eu-west1) — MCI handles N clusters natively"],
]
_table_header_row(tbl, risks[0])
for r in range(1, 6):
    _table_data_row(tbl, r, risks[r])

add_notes(s, """MEETING NOTES — Risks
- SSL cert provisioning is the most common 'stuck' issue — always verify DNS first
- MCI API enablement handled in Terraform but can fail if org policies block it
- Binary Authorization: recommend dry-run first day to catch image signing issues
- BigQuery cost: pre-aggregation + filter DEBUG/TRACE logs at the sink
- Region expansion: add 3rd gke-cluster module call and register with fleet
- No showstoppers identified — all risks have known mitigations""")


# ── Slide 15: Deliverables ────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Deliverables Summary")
add_footer(s)

rows, cols = 9, 2
tbl_shape = s.shapes.add_table(rows, cols,
    Inches(0.6), CTOP, Inches(12.1), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(5.8)
tbl.columns[1].width = Inches(6.3)

deliverables = [
    ["Deliverable",                           "Location"],
    ["Terraform (entire infrastructure)",     "terraform/"],
    ["Architecture diagrams (Mermaid + SVG)", "docs/"],
    ["Step-by-step setup instructions",       "README.md, terraform/README.md"],
    ["BigQuery schema + sample queries",      "observability/bigquery/"],
    ["Design decisions & rationale",          "docs/design-decisions.md"],
    ["Sample apps + Dockerfiles",             "apps/"],
    ["CI/CD pipelines",                       "cicd/"],
    ["Grafana dashboard + datasource",        "observability/grafana/"],
]
_table_header_row(tbl, deliverables[0])
for r in range(1, 9):
    _table_data_row(tbl, r, deliverables[r])

add_notes(s, """MEETING NOTES — Deliverables
- Everything in a single Git repository — version-controlled and auditable
- Terraform modules are reusable across environments (dev, prod)
- Architecture diagrams are Mermaid format (text-based, diffable) with SVG exports
- Runbook covers day-2 operations: scaling, failover, troubleshooting
- BigQuery sample queries help ops teams investigate incidents without full schema knowledge
- All deliverables listed are complete and committed to the repository""")


# ── Slide 16: Next Steps & Ask ────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s, "Next Steps & Ask")
add_footer(s)
add_body(s, [
    ("Immediate Next Steps", 0),
    ("Secure GCP org/billing approval for the new project", 1),
    ("Assign 2-person engineering team (1 platform, 1 app/observability)", 1),
    ("Confirm DNS domain and hostname strategy", 1),
    ("Set up GitHub repository with branch protection + required reviews", 1),
    ("", 0),
    ("Ask from Leadership", 0),
    ("Approval to proceed with the 2-week sprint plan", 1),
    ("Confirm environment scope: Dev only first, or Dev + Prod in parallel?", 1),
    ("Decision on Anthos Service Mesh (mTLS) — adds ~3 days if included", 1),
    ("Budget approval for GCP resources (~$X/month estimated)", 1),
    ("", 0),
    ("Out of Scope (unless explicitly added)", 0),
    ("Anthos Service Mesh / mTLS  |  Shared VPC  |  Hybrid/on-prem fleet  |  Synthetic monitoring", 1),
], font_size=16)
add_notes(s, """MEETING NOTES — Next Steps
- Critical path blocker is org/billing approval — start immediately
- 2-person team is the sweet spot: fast (10–12 days) without coordination overhead
- DNS decision needed: what domain? Subdomains? Existing DNS zones?
- Dev environment first is recommended — validate everything, then clone to prod with different tfvars
- Anthos Service Mesh adds mTLS between services and traffic splitting — valuable but +3 days
- After approval, Day 1 starts with bootstrap.sh to create state bucket and WIF""")


# ═══════════════════════════════════════════════════════════════════════════════
# CLOSING SLIDE  (Slide 17)
# ═══════════════════════════════════════════════════════════════════════════════
# Mirror the title slide layout — all shapes, no images.

slide17 = prs.slides.add_slide(prs.slide_layouts[6])

# Full dark-blue background
_rect(slide17, 0, 0, W, prs.slide_height, C_DBLUE)

# Dark navy banner area
_rect(slide17, 0, 0, W, TITLE_BANNER_H, C_NAVY)

# Branding text inside banner
_textbox(slide17, 0.6, 0.22, 6, 0.35,
         "●  PRODUCTION-GRADE · GA", 11, C_LBLUE)
_textbox(slide17, 0.6, 0.58, 9, 1.10,
         "GCP × SRE × O11y", 52, C_WHITE, bold=True)
_textbox(slide17, 0.6, 1.58, 11, 0.45,
         "Multi-region GKE platform with Grafana-on-BigQuery observability",
         17, C_LBLUE)

# Feature pills
for lbl, px in zip(pill_labels, pill_x):
    pill = _rect(slide17, Inches(px), Inches(1.90), Inches(2.10), Inches(0.30), C_DBLUE)
    pill.line.color.rgb = C_GBLUE
    tf_p = pill.text_frame
    tf_p.word_wrap = False
    tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_p.margin_left = tf_p.margin_right = Inches(0.08)
    tf_p.margin_top = tf_p.margin_bottom = 0
    pp = tf_p.paragraphs[0]
    pp.text = lbl
    pp.font.size = Pt(11)
    pp.font.color.rgb = C_VLBLUE
    pp.alignment = PP_ALIGN.CENTER

# Accent strip
for i, c in enumerate(GOOGLE4):
    _rect(slide17, seg_w * i, TITLE_BANNER_H, seg_w, TITLE_STRIP_H, c)

# Thank you text
_textbox(slide17, 0.8, 2.65, 11.7, 1.00,
         "Thank You", 54, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
_textbox(slide17, 0.8, 3.70, 11.7, 0.55,
         "Questions & Discussion",
         26, C_VLBLUE, align=PP_ALIGN.CENTER)
_textbox(slide17, 0.8, 4.55, 11.7, 0.55,
         "AI Engineer / SRE Architect — Abhay Ojha",
         16, C_LBLUE, align=PP_ALIGN.CENTER)

add_footer(slide17, on_dark=True)

add_notes(slide17, """MEETING NOTES — Closing
- Open the floor for questions
- Key decisions needed from leadership:
  1. Approval to proceed
  2. Team assignment
  3. Environment scope (dev first vs dev+prod)
  4. Anthos Service Mesh yes/no
  5. Budget sign-off
- Offer to schedule a follow-up deep-dive on any specific area
- Share the repository link and docs for async review""")


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
OUTPUT = "/Users/abojha/Anthropic-POC/GCP-Grafana/GCP - Grafana/GCP_Platform_Leadership_Review.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Header layout: branding={BRAND_H/914400:.3f}\"  strip={STRIP_H/914400:.3f}\"  title={TITLB_H/914400:.3f}\"  total={HDR_H/914400:.3f}\"")
print(f"Content starts at: {CTOP/914400:.3f}\"  Body height: {BODY_H/914400:.3f}\"")
